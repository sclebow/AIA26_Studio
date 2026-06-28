"""Build an editable PowerPoint (.pptx) of the TerraPilot presentation.

Mirrors the content of terrapilot-deck.html as native, EDITABLE PowerPoint slides
(text boxes, bullet lists, tables, coloured cards) on a dark TerraPilot theme. The
HTML deck's live animations (state graph, shape-morph, Pareto scatter) become clean
static diagrams here — open the .pptx in PowerPoint / Google Slides and edit freely.

Run:  python build_pptx.py   ->  writes TerraPilot.pptx next to this file.
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# ---- TerraPilot theme (matches the HTML deck) ----
BG      = RGBColor(0x07, 0x0B, 0x14)   # deep navy
PANEL   = RGBColor(0x0D, 0x16, 0x28)   # card fill
PANEL2  = RGBColor(0x10, 0x1B, 0x30)
INK     = RGBColor(0xEA, 0xF2, 0xFF)   # bright text
MUTED   = RGBColor(0x8B, 0x97, 0xAD)   # dim text
CYAN    = RGBColor(0x28, 0xE0, 0xD0)
CYAN2   = RGBColor(0x57, 0xF2, 0xE6)
VIOLET  = RGBColor(0x7C, 0x7B, 0xFF)
AMBER   = RGBColor(0xFF, 0xB4, 0x54)
PINK    = RGBColor(0xFF, 0x8F, 0xA3)
GREEN   = RGBColor(0x57, 0xE0, 0x8A)
LINE    = RGBColor(0x2A, 0x38, 0x52)

# 16:9 canvas
EMU = 914400
SW, SH = int(13.333 * EMU), int(7.5 * EMU)
FONT = "Segoe UI"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def inch(v):
    return Emu(int(v * EMU))


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)   # rectangle
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    # subtle top accent bar
    bar = s.shapes.add_shape(1, 0, 0, SW, inch(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background(); bar.shadow.inherit = False
    return s


def _set(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = font


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = inch(0.05)
    tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, italic=False, first=False,
         align=PP_ALIGN.LEFT, font=FONT, space_after=4, bullet=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
    if bullet:
        text = "▸  " + text
    r = p.add_run(); r.text = text
    _set(r, size, color, bold, italic, font)
    return p


def kicker(s, text):
    tf = textbox(s, 0.7, 0.45, 11, 0.5)
    para(tf, text.upper(), 13, CYAN, bold=True, first=True)


def title(s, plain, hl=None, tail=""):
    tf = textbox(s, 0.7, 0.85, 12, 1.0)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = plain; _set(r, 33, INK, bold=True)
    if hl:
        r2 = p.add_run(); r2.text = hl; _set(r2, 33, CYAN2, bold=True)
    if tail:
        r3 = p.add_run(); r3.text = tail; _set(r3, 33, INK, bold=True)


def card(s, x, y, w, h, icon, head, body, accent=CYAN):
    box = s.shapes.add_shape(1, inch(x), inch(y), inch(w), inch(h))
    box.fill.solid(); box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = LINE; box.line.width = Pt(1); box.shadow.inherit = False
    # left accent strip
    strip = s.shapes.add_shape(1, inch(x), inch(y), inch(0.06), inch(h))
    strip.fill.solid(); strip.fill.fore_color.rgb = accent
    strip.line.fill.background(); strip.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = inch(0.2); tf.margin_right = inch(0.15)
    tf.margin_top = inch(0.15); tf.margin_bottom = inch(0.1)
    para(tf, f"{icon}  {head}", 15, accent, bold=True, first=True, space_after=5)
    para(tf, body, 11.5, RGBColor(0xC7, 0xD2, 0xE6))


def speaker(s, text):
    tf = textbox(s, 0.7, 6.75, 12, 0.55)
    p = para(tf, text, 11, VIOLET, italic=True, first=True)


def panel(s, x, y, w, h, fill=PANEL, accent=LINE):
    box = s.shapes.add_shape(1, inch(x), inch(y), inch(w), inch(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = accent; box.line.width = Pt(1.25); box.shadow.inherit = False
    return box


def bullets(s, x, y, w, h, items, size=13.5):
    """items: list of (bold_lead, small_sub) tuples."""
    tf = textbox(s, x, y, w, h)
    first = True
    for lead, sub in items:
        p = para(tf, lead, size, INK, bold=True, first=first, bullet=True, space_after=1)
        first = False
        if sub:
            sp = tf.add_paragraph(); sp.space_after = Pt(7); sp.space_before = Pt(0)
            r = sp.add_run(); r.text = "     " + sub; _set(r, 10.5, MUTED)


# ====================== SLIDE 1 — TITLE ======================
s = slide()
kicker(s, "AI-Driven Architectural Design")
tf = textbox(s, 0.7, 2.2, 12, 1.4)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "▲ Terra"; _set(r, 60, INK, bold=True)
r2 = p.add_run(); r2.text = "Pilot"; _set(r2, 60, CYAN2, bold=True)
tf2 = textbox(s, 0.7, 3.6, 11.5, 1.0)
para(tf2, "An AI co-pilot that turns plain prompts into site-aware, optimized, "
          "explainable architecture.", 20, MUTED, first=True)
# team strip
labels = ["AI Logic & Optimization", "Shape & Manipulation", "Site & Urban Context", "Overview"]
cw = 2.9
for i, lab in enumerate(labels):
    b = panel(s, 0.7 + i * (cw + 0.12), 5.0, cw, 0.7, fill=PANEL)
    tf3 = b.text_frame; tf3.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf3, lab, 12, CYAN2, bold=True, first=True, align=PP_ALIGN.CENTER)
speaker(s, '"You describe a building in a simple prompt — TerraPilot generates, places, and optimizes it on a real site."')

# ====================== SLIDE 2 — THE IDEA ======================
s = slide()
kicker(s, "The Idea")
title(s, "From a ", "sentence", " to a site-aware building")
tf = textbox(s, 0.7, 1.7, 12, 0.8)
para(tf, "Early-stage design is slow and manual — one massing study can take days. "
         "TerraPilot collapses that into minutes through a guided, conversational workflow.",
     14, MUTED, first=True)
cards = [
    ("💬", "You talk to it", "No menus to learn. Just say it — “give it a futuristic twist”, "
     "“open it toward the park” — like talking to a design assistant.", CYAN),
    ("🌍", "On a real site", "It knows the actual roads, sun, wind and neighbours around your "
     "plot — not a blank canvas.", VIOLET),
    ("🧠", "And it explains", "Every change carries an intent, a confidence, and a WHY. "
     "It tells you what's real vs estimated.", CYAN2),
]
for i, (ic, h, b, a) in enumerate(cards):
    card(s, 0.7 + i * 4.05, 2.7, 3.8, 1.9, ic, h, b, a)
# stat row
stats = [("2000 m", "OSM context radius"), ("16", "scoring objectives"),
         ("8", "named strategies"), ("3", "understanding layers")]
for i, (n, c) in enumerate(stats):
    b = panel(s, 0.7 + i * 3.05, 4.85, 2.9, 1.05)
    tf2 = b.text_frame; tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf2, n, 26, CYAN2, bold=True, first=True, align=PP_ALIGN.CENTER, space_after=2)
    para(tf2, c, 11, MUTED, align=PP_ALIGN.CENTER)

# ====================== SLIDE 3 — WORKFLOW ======================
s = slide()
kicker(s, "The Journey")
title(s, "One guided ", "workflow", ", seven stages")
stages = [("📍", "01 Site", "Define the plot."), ("⬡", "02 Boundary", "Auto setbacks."),
          ("🛰", "03 Context", "OSM digital twin."), ("◈", "04 Shape", "Generate & edit."),
          ("⚙", "05 Optimize", "Multi-objective Pareto."), ("▦", "06 Compare", "Side by side."),
          ("⬇", "07 Export", "Rhino / Revit IFC.")]
for i, (ic, h, b) in enumerate(stages):
    col = i % 4; rowi = i // 4
    x = 0.7 + col * 3.05; y = 2.1 + rowi * 1.7
    bx = panel(s, x, y, 2.9, 1.45, fill=PANEL, accent=LINE)
    tf2 = bx.text_frame; tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf2, ic, 20, CYAN, first=True, align=PP_ALIGN.CENTER, space_after=2)
    para(tf2, h, 14, INK, bold=True, align=PP_ALIGN.CENTER, space_after=1)
    para(tf2, b, 11, MUTED, align=PP_ALIGN.CENTER)
speaker(s, "Backend enforces the order: shape generation + selection must precede optimization.")

# ====================== SLIDE 4 — ARCHITECTURE ======================
s = slide()
kicker(s, "System Design")
title(s, "How the ", "frontend & backend", " work together")
fe = panel(s, 0.7, 1.8, 5.7, 4.6, accent=CYAN)
tf = fe.text_frame; tf.margin_left = inch(0.25); tf.margin_top = inch(0.2); tf.word_wrap = True
para(tf, "🖥  Frontend · vanilla JS", 16, CYAN2, bold=True, first=True, space_after=8)
for code, sub in [("copilotClient.js", "Routes your message to the right action"),
                  ("core/store.js", "Single source of truth — UI reacts to state"),
                  ("views/viewerView.js", "3D building — extrudes each floor plate at its Z"),
                  ("panels/healthDashboard.js", "Live scores, re-evaluates on every edit")]:
    para(tf, code, 13, CYAN, bold=True, font=MONO, space_after=0)
    para(tf, "   " + sub, 10.5, MUTED, space_after=7)
be = panel(s, 6.9, 1.8, 5.7, 4.6, accent=VIOLET)
tf = be.text_frame; tf.margin_left = inch(0.25); tf.margin_top = inch(0.2); tf.word_wrap = True
para(tf, "⚙  Backend · FastAPI", 16, RGBColor(0xC9, 0xC8, 0xFF), bold=True, first=True, space_after=8)
for code, sub in [("connection/routes/", "API endpoints — thin wrappers, no duplicated logic"),
                  ("connection/nodes/", "3-layer prompt understanding"),
                  ("notebook_logic/", "Geometry engine — manipulate, transform, floor plates"),
                  ("optimization/", "16 scorers + Pareto selector + env data")]:
    para(tf, code, 13, VIOLET, bold=True, font=MONO, space_after=0)
    para(tf, "   " + sub, 10.5, MUTED, space_after=7)
speaker(s, "One shared runtime — the notebooks and the UI call the same engine. JSON over HTTP between them.")

# ====================== SLIDE 5 — PROMPT FLOW ======================
s = slide()
kicker(s, "End to End")
title(s, "Watch a ", "prompt flow", " through the system")
tf = textbox(s, 0.7, 1.7, 12, 0.5)
para(tf, "“Give it a futuristic twist” — from the chat box to the sculpted 3D building.",
     14, MUTED, italic=True, first=True)
steps = [("💬", "Chat input", "User types intent", "copilotClient.js"),
         ("🔀", "API route", "/feedback endpoint", "feedback_routes.py"),
         ("🧠", "Understand", "3-layer → intent + WHY", "nodes/"),
         ("📐", "Geometry op", "twist the plate stack", "notebook_logic"),
         ("🧊", "3D viewer", "building updates live", "viewerView.js")]
bw = 2.3
for i, (ic, h, sub, tag) in enumerate(steps):
    x = 0.55 + i * 2.5
    bx = panel(s, x, 2.5, bw, 1.7, fill=PANEL, accent=LINE)
    tf2 = bx.text_frame; tf2.vertical_anchor = MSO_ANCHOR.MIDDLE; tf2.word_wrap = True
    para(tf2, ic, 19, CYAN, first=True, align=PP_ALIGN.CENTER, space_after=2)
    para(tf2, h, 13, INK, bold=True, align=PP_ALIGN.CENTER, space_after=1)
    para(tf2, sub, 10, MUTED, align=PP_ALIGN.CENTER, space_after=3)
    para(tf2, tag, 9.5, CYAN2, font=MONO, align=PP_ALIGN.CENTER)
    if i < 4:
        ar = textbox(s, x + bw - 0.05, 3.0, 0.6, 0.6, anchor=MSO_ANCHOR.MIDDLE)
        para(ar, "→", 22, VIOLET, first=True, align=PP_ALIGN.CENTER)
foot = ["① Intent: change_character", "② Op: twist · confidence 0.9",
        "③ Plates re-transformed, holes kept", "④ Response → store → viewer re-renders"]
tf = textbox(s, 0.7, 4.6, 12, 1.0)
para(tf, "      ".join(foot), 12, CYAN2, first=True)
speaker(s, 'Every hop is real — the same path runs for a command, a vague intent, or "align to the road".')

# ====================== SLIDE 6 — DIGITAL TWIN ======================
s = slide()
kicker(s, "Stage 1–3")
title(s, "A ", "digital twin", " of the surroundings")
bullets(s, 0.7, 1.9, 6.6, 4.2, [
    ("Real OpenStreetMap data in a 2 km radius", "Roads, transit, parks, amenities & neighbouring buildings"),
    ("Edge intelligence — every site edge named by compass", '"north edge", "road-facing edge" become design targets'),
    ("10 context scores computed once & cached", "Feeds noise, view, access & density objectives later"),
    ("Automatic setbacks define the buildable area", "Nothing can be placed outside it"),
])
# live dashboard panel
dp = panel(s, 7.6, 1.9, 5.0, 4.1, accent=CYAN)
tf = dp.text_frame; tf.margin_left = inch(0.25); tf.margin_top = inch(0.2); tf.word_wrap = True
para(tf, "● Live Design Health", 15, CYAN2, bold=True, first=True, space_after=8)
for lbl, val, col in [("☀ Solar", 66, AMBER), ("👁 View", 98, CYAN2), ("🔊 Noise", 14, VIOLET),
                      ("💨 Wind", 64, RGBColor(0x38, 0xBD, 0xF8)), ("🏢 Density", 56, PINK)]:
    pr = tf.add_paragraph(); pr.space_after = Pt(8)
    r = pr.add_run(); r.text = f"{lbl}"; _set(r, 12.5, INK)
    r2 = pr.add_run(); r2.text = f"   {'█' * (val // 7)}  {val}"; _set(r2, 12.5, col, bold=True)
para(tf, "Draggable · minimizable · updates live · shows geometry → real data", 10, MUTED, space_after=0)

# ====================== SLIDE 7 — MANIPULATION ======================
s = slide()
kicker(s, "Stage 4 · Manipulation")
title(s, "Sculpt the building by ", "talking", " to it")
panel(s, 0.7, 1.9, 5.6, 4.2, accent=VIOLET)
tf = textbox(s, 0.95, 2.2, 5.1, 3.6, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "💬  “Design me an 8-storey tower”", 16, CYAN2, bold=True, first=True, space_after=10)
para(tf, "→ “give it a futuristic twist”", 14, INK, space_after=6)
para(tf, "→ “carve a courtyard”", 14, INK, space_after=6)
para(tf, "→ “add 2 floors on the right wing”", 14, INK, space_after=10)
para(tf, "Real geometry updates live in 3D after every prompt.", 11.5, MUTED, italic=True)
bullets(s, 6.7, 2.0, 5.9, 4.0, [
    ("Natural-language intent, not commands", "“Give it a futuristic twist”, “carve a courtyard”"),
    ("Floor-plate model — a stack of plates", "Edit per-floor AND per-wing, not the whole block"),
    ("Real geometry every time", "Holes, wings & twist travel through every transform"),
    ("Save any version", "→ preview & optimize it later"),
])

# ====================== SLIDE 8 — PROMPT UNDERSTANDING ======================
s = slide()
kicker(s, "The Brain")
title(s, "How we ", "break & understand", " a prompt")
tf = textbox(s, 0.7, 1.7, 12, 0.5)
para(tf, "Every message flows through a 3-layer pipeline — most specific first, smartest fallback last.",
     14, MUTED, first=True)
layers = [("⌨️", "LAYER 1", "Direct command", "Exact instructions — “add 2 floors”, “rotate 30°”.", CYAN),
          ("🧩", "LAYER 2", "Semantic intent", "The LLM reads a feeling — “let in more daylight”.", VIOLET),
          ("🧠", "LAYER 3", "Reason Node", "Decides the move + the WHY, with confidence.", CYAN2)]
for i, (ic, n, h, b, a) in enumerate(layers):
    x = 0.7 + i * 4.05
    bx = panel(s, x, 2.35, 3.8, 1.85, accent=a)
    tf2 = bx.text_frame; tf2.margin_left = inch(0.2); tf2.margin_top = inch(0.15); tf2.word_wrap = True
    para(tf2, f"{ic}  {n}", 12, a, bold=True, first=True, space_after=3)
    para(tf2, h, 15, INK, bold=True, space_after=4)
    para(tf2, b, 11, MUTED)
    if i < 2:
        ar = textbox(s, x + 3.8, 2.9, 0.45, 0.6, anchor=MSO_ANCHOR.MIDDLE)
        para(ar, "→", 22, VIOLET, first=True, align=PP_ALIGN.CENTER)
card(s, 0.7, 4.5, 5.8, 1.6, "🔁", "Two-tier fallback",
     "Gemini → Cloudflare → keyword matching. It never dies on an API quota — the fallback is load-bearing.", CYAN)
card(s, 6.8, 4.5, 5.8, 1.6, "🛡️", "Hallucination guard",
     "If a prompt matches no real operation, we reject it — never fake a change to look successful.", PINK)

# ====================== SLIDE 9 — OPTIMIZATION ======================
s = slide()
kicker(s, "The Engine")
title(s, "Multi-objective ", "optimization")
steps = [("🎲 POOL", "Candidates", "Placement × rotation × massing — large, diverse pool."),
         ("📊 SCORE", "Every objective", "16 scorers per candidate, cached site memory."),
         ("⭐ SELECT", "8 strategies", "Diverse, named options from the Pareto front.")]
for i, (n, h, b) in enumerate(steps):
    bx = panel(s, 0.7, 2.0 + i * 1.45, 5.7, 1.3, accent=CYAN)
    tf2 = bx.text_frame; tf2.margin_left = inch(0.2); tf2.vertical_anchor = MSO_ANCHOR.MIDDLE; tf2.word_wrap = True
    para(tf2, f"{n} — {h}", 14, CYAN2, bold=True, first=True, space_after=2)
    para(tf2, b, 11, MUTED)
rp = panel(s, 6.9, 2.0, 5.7, 4.0, accent=VIOLET)
tf = rp.text_frame; tf.margin_left = inch(0.25); tf.margin_top = inch(0.2); tf.word_wrap = True
para(tf, "True Pareto front", 15, CYAN2, bold=True, first=True, space_after=6)
para(tf, "A diverse set of non-dominated options — each best at a different trade-off, "
         "so you choose by priority, not a single 'winner'.", 12, INK, space_after=10)
para(tf, "☀ Solar    👁 View    🔊 Noise", 12, CYAN, bold=True, space_after=3)
para(tf, "🌳 Open    💨 Wind    ⚖ Balanced", 12, CYAN, bold=True)
speaker(s, "Honest engineering: scores use real NASA POWER climate + OSM data. Missing data shows N/A, never invented.")

# ====================== SLIDE 10 — HARD PROBLEMS ======================
s = slide()
kicker(s, "Depth")
title(s, "The hard problems we ", "solved")
cards3 = [
    ("🌀", "Optimize keeps your edits", "Optimizing a manipulated building preserves the twist + "
     "added floors — the full floor-plate stack is carried through every step.", CYAN),
    ("🗂️", "Saved-shapes explorer", "Browse saved designs → preview each in 3D → optimize the "
     "exact one you choose. Variants accumulate, never overwrite.", VIOLET),
    ("🛣️", "Context-aware alignment", '"Align facade to main road" rotates the building parallel '
     "to the real OSM road geometry — not a compass default.", CYAN2),
]
for i, (ic, h, b, a) in enumerate(cards3):
    card(s, 0.7 + i * 4.05, 2.0, 3.8, 2.5, ic, h, b, a)
pb = panel(s, 0.7, 4.8, 11.9, 1.3, accent=GREEN)
tf = pb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = inch(0.25); tf.word_wrap = True
para(tf, "All 6 environmental scorers run in dual mode — a geometry estimate AND a real-data "
         "value — with a proxy / ● live badge, so you always know which is which.",
     14, INK, bold=True, first=True)

# ====================== SLIDE 11 — BACKEND ======================
s = slide()
kicker(s, "Under the Hood")
title(s, "The ", "backend", ", in one picture")
fp = panel(s, 0.7, 2.0, 5.7, 3.7, accent=CYAN)
tf = fp.text_frame; tf.margin_left = inch(0.25); tf.margin_top = inch(0.2); tf.word_wrap = True
para(tf, "FastAPI runtime", 15, CYAN2, bold=True, first=True, space_after=8)
for lead, sub in [("One shared engine — notebooks & UI", "Routes are thin wrappers over notebook_logic"),
                  ("nodes/ — 3-layer prompt understanding", "direct → semantic → reason-node"),
                  ("optimization/ — 16 scorers + Pareto", "env_data · site_context_memory · context_enrichment")]:
    para(tf, "▸ " + lead, 13, INK, bold=True, space_after=0)
    para(tf, "   " + sub, 10.5, MUTED, space_after=7)
dp = panel(s, 6.9, 2.0, 5.7, 3.7, accent=VIOLET)
tf = dp.text_frame; tf.margin_left = inch(0.25); tf.margin_top = inch(0.2); tf.word_wrap = True
para(tf, "Data & intelligence", 15, RGBColor(0xC9, 0xC8, 0xFF), bold=True, first=True, space_after=8)
for lead, sub in [("NASA POWER — wind, solar, temperature", "Fetched once per site, cached"),
                  ("OpenStreetMap — roads, amenities, buildings", "Powers noise, access, density, road-align"),
                  ("Shape version manager", "Save / preview / optimize — floor-plate geometry preserved")]:
    para(tf, "▸ " + lead, 13, INK, bold=True, space_after=0)
    para(tf, "   " + sub, 10.5, MUTED, space_after=7)

# ====================== SLIDE 12 — AGENT STATE GRAPH ======================
s = slide()
kicker(s, "Under the Hood")
title(s, "The agent's ", "state graph", " — every prompt's path")
tf = textbox(s, 0.7, 1.7, 12, 0.55)
para(tf, "A supervisor (central_reason) picks the next action; each node writes its result to "
         "shared state and returns to the planner. Drawn straight from agent/graph.py.",
     13, MUTED, first=True)
# simple node ladder
def node_box(x, y, w, h, label, sub, col):
    bx = panel(s, x, y, w, h, accent=col)
    tf2 = bx.text_frame; tf2.vertical_anchor = MSO_ANCHOR.MIDDLE; tf2.word_wrap = True
    para(tf2, label, 12, col, bold=True, first=True, align=PP_ALIGN.CENTER, space_after=1, font=MONO)
    if sub:
        para(tf2, sub, 9.5, MUTED, align=PP_ALIGN.CENTER)
node_box(5.4, 2.35, 2.5, 0.55, "START → planner", "builds the plan", MUTED)
node_box(5.0, 3.05, 3.3, 0.7, "central_reason", "supervisor — picks next action", PINK)
actions = [("read_site", CYAN), ("generate_shape", VIOLET), ("check_constraints", PINK),
           ("optimize", GREEN), ("evaluate", CYAN2), ("place_building", VIOLET),
           ("report → finish", GREEN), ("await_human → finish", RGBColor(0x5E, 0xC8, 0xFF))]
for i, (lbl, col) in enumerate(actions):
    x = 0.7 + (i % 4) * 3.05; y = 4.05 + (i // 4) * 0.85
    node_box(x, y, 2.9, 0.65, lbl, "↑ writes to state → planner", col)
tf = textbox(s, 0.7, 5.95, 12, 0.7)
para(tf, "Nodes don't hand off to each other — each writes its result to shared state, then returns "
         "to the planner, which reads it and plans the next step.", 12, RGBColor(0x9F, 0xB0, 0xC8),
     bold=True, first=True)

# ====================== SLIDE 13 — WHAT WORKS / LIMITS ======================
s = slide()
kicker(s, "Honest Status")
title(s, "What ", "works", " — and what doesn't yet")
ok = panel(s, 0.7, 2.0, 5.85, 4.1, accent=GREEN)
tf = ok.text_frame; tf.margin_left = inch(0.22); tf.margin_top = inch(0.18); tf.word_wrap = True
para(tf, "✅ What works today", 14, GREEN, bold=True, first=True, space_after=6)
for t in ["Plain-prompt generation, manipulation & per-wing floor edits",
          "Multi-objective optimization with a true Pareto front",
          "Real environmental scoring — NASA + OpenStreetMap",
          'Context-aware moves — "align to the road", "face the park"',
          "Save / preview / optimize saved versions end-to-end"]:
    para(tf, "• " + t, 11.5, INK, space_after=5)
lim = panel(s, 6.75, 2.0, 5.85, 4.1, accent=AMBER)
tf = lim.text_frame; tf.margin_left = inch(0.22); tf.margin_top = inch(0.18); tf.word_wrap = True
para(tf, "🚧 Current limitations", 14, AMBER, bold=True, first=True, space_after=6)
for t in ["A scoring engine, not a full physics simulation",
          "LLM intent depends on a free-tier quota (keyword fallback covers it)",
          "Export is GeoJSON today; native Rhino/Revit needs the MCP bridge",
          "Single building per site; massing is conceptual",
          "Context quality depends on OpenStreetMap coverage"]:
    para(tf, "• " + t, 11.5, INK, space_after=5)
speaker(s, "We're deliberate about honesty: when data or a model isn't available, TerraPilot says so — never invents a number.")

# ====================== SLIDE 14 — VALUE TO DESIGNERS ======================
s = slide()
kicker(s, "Why designers care")
title(s, "How it ", "helps", " a designer")
vals = [("⚡", "Explore in minutes, not days", "Test dozens of site-aware options in the time a single massing study used to take."),
        ("🌍", "Site-grounded from the first sketch", "Sun, wind, noise, roads and neighbours are in the loop from minute one."),
        ("🗣️", "No software to learn", "Designers think in intent, not menus. “Open it toward the park” just works."),
        ("🧾", "Decisions you can defend", "Every move carries a WHY and a score trade-off — a clear story for clients & juries.")]
for i, (ic, h, b) in enumerate(vals):
    x = 0.7 + (i % 2) * 6.05; y = 2.0 + (i // 2) * 2.05
    card(s, x, y, 5.8, 1.85, ic, h, b, CYAN if i % 2 == 0 else VIOLET)
speaker(s, "TerraPilot doesn't replace the designer — it removes the grunt work so they spend time on judgement, not setup.")

# ====================== SLIDE 15 — CLOSE ======================
s = slide()
tf = textbox(s, 0.7, 2.6, 12, 1.6, anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Plain prompts in.\n"; _set(r, 40, INK, bold=True)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "Site-aware architecture out."; _set(r2, 40, CYAN2, bold=True)
tf2 = textbox(s, 0.7, 4.6, 12, 0.8, anchor=MSO_ANCHOR.MIDDLE)
para(tf2, "▲ TerraPilot · An AI co-pilot for early-stage architectural design",
     18, MUTED, first=True, align=PP_ALIGN.CENTER)

# ---- save ----
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "TerraPilot.pptx")
prs.save(out)
print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
